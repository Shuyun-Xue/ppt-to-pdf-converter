import streamlit as st
import os
import tempfile
from pptx import Presentation
from pptx.util import Inches
from fpdf import FPDF
from PIL import Image, ImageDraw
import io
from PyPDF2 import PdfReader, PdfWriter
import hashlib
import time

# 常量定义
MAX_FILE_SIZE = 200 * 1024 * 1024  # 200MB
CACHE_DIR = "cache"
CHUNK_SIZE = 8 * 1024 * 1024  # 8MB 分块大小

# 创建缓存目录
os.makedirs(CACHE_DIR, exist_ok=True)

@st.cache_data
def get_file_hash(file_bytes):
    """计算文件哈希值用于缓存"""
    return hashlib.md5(file_bytes).hexdigest()

def format_size(size_bytes):
    """格式化文件大小显示"""
    for unit in ['B', 'KB', 'MB', 'GB']:
        if size_bytes < 1024.0:
            return f"{size_bytes:.1f} {unit}"
        size_bytes /= 1024.0
    return f"{size_bytes:.1f} GB"

def save_uploaded_file(uploaded_file, target_path, progress_bar=None):
    """分块保存上传的文件，显示进度"""
    file_size = len(uploaded_file.getvalue())
    with open(target_path, 'wb') as f:
        bytes_data = uploaded_file.getvalue()
        total_chunks = (file_size + CHUNK_SIZE - 1) // CHUNK_SIZE
        
        for i in range(0, file_size, CHUNK_SIZE):
            chunk = bytes_data[i:i + CHUNK_SIZE]
            f.write(chunk)
            if progress_bar:
                progress_bar.progress((i + len(chunk)) / file_size)

def compress_pdf(input_path, quality='medium', progress_bar=None):
    """压缩PDF文件
    quality: 'low', 'medium', 'high'
    """
    try:
        # 创建临时文件来保存压缩后的PDF
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_file:
            output_path = tmp_file.name

        # 读取PDF
        reader = PdfReader(input_path)
        writer = PdfWriter()

        # 压缩质量设置
        compression_params = {
            'low': {'image_quality': 30, 'compress_pages': True},
            'medium': {'image_quality': 50, 'compress_pages': True},
            'high': {'image_quality': 70, 'compress_pages': True}
        }
        params = compression_params[quality]

        # 处理每一页
        total_pages = len(reader.pages)
        for i, page in enumerate(reader.pages):
            writer.add_page(page)
            writer.add_metadata(reader.metadata)
            if progress_bar:
                progress_bar.progress((i + 1) / (total_pages + 1))
            
        # 保存压缩后的文件
        with open(output_path, 'wb') as output_file:
            writer.write(output_file)

        return output_path
    except Exception as e:
        st.error(f"PDF压缩出错: {str(e)}")
        return None

def render_shape(shape, draw, offset_x=0, offset_y=0):
    """渲染PPT中的形状"""
    try:
        if hasattr(shape, 'text'):
            # 渲染文本
            text_frame = shape.text_frame
            if text_frame.text:
                draw.text((offset_x + shape.left, offset_y + shape.top), 
                         text_frame.text, 
                         fill='black')
        
        if hasattr(shape, 'fill'):
            # 渲染形状
            if shape.shape_type == 1:  # 矩形
                draw.rectangle(
                    [offset_x + shape.left, offset_y + shape.top,
                     offset_x + shape.left + shape.width,
                     offset_y + shape.top + shape.height],
                    outline='black'
                )
    except Exception as e:
        st.warning(f"渲染形状时出现警告: {str(e)}")

def convert_slide_to_image(slide, progress_bar=None):
    """将PPT幻灯片转换为图像"""
    # 获取幻灯片尺寸
    width = int(Inches(slide.slide_width).px)
    height = int(Inches(slide.slide_height).px)
    
    # 创建一个新的图像
    img = Image.new('RGB', (width, height), 'white')
    draw = ImageDraw.Draw(img)
    
    # 渲染所有形状
    total_shapes = len(slide.shapes)
    for i, shape in enumerate(slide.shapes):
        render_shape(shape, draw)
        if progress_bar:
            progress_bar.progress((i + 1) / total_shapes)
    
    return img

def convert_ppt_to_pdf(input_file_path, compression_quality='medium'):
    """
    将PPT转换为PDF（纯Python实现）
    compression_quality: 'low', 'medium', 'high'
    """
    try:
        # 创建临时目录存放图片
        with tempfile.TemporaryDirectory() as temp_dir:
            # 加载PPT文件
            prs = Presentation(input_file_path)
            
            # 创建PDF文档
            pdf = FPDF()
            
            # 设置PDF页面大小为PPT大小
            first_slide = prs.slides[0] if prs.slides else None
            if first_slide:
                width = Inches(prs.slide_width).inches * 25.4  # 转换为毫米
                height = Inches(prs.slide_height).inches * 25.4
                pdf.set_page_size((width, height))
            
            # 创建进度条
            progress_text = "转换进度"
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            # 遍历所有幻灯片
            total_slides = len(prs.slides)
            for i, slide in enumerate(prs.slides):
                status_text.text(f"正在处理第 {i+1}/{total_slides} 页...")
                
                # 将幻灯片转换为图像
                img = convert_slide_to_image(slide)
                
                # 保存图像
                img_path = os.path.join(temp_dir, f'slide_{i}.png')
                img.save(img_path, 'PNG', optimize=True)
                
                # 添加到PDF
                pdf.add_page()
                pdf.image(img_path, x=0, y=0, w=pdf.w, h=pdf.h)
                
                # 更新进度
                progress_bar.progress((i + 1) / total_slides)
            
            # 保存PDF
            status_text.text("正在生成PDF...")
            output_path = os.path.splitext(input_file_path)[0] + '.pdf'
            pdf.output(output_path)
            
            # 压缩PDF
            if compression_quality != 'none':
                status_text.text("正在压缩PDF...")
                compressed_path = compress_pdf(output_path, compression_quality, progress_bar)
                if compressed_path:
                    # 删除原始PDF
                    os.remove(output_path)
                    output_path = compressed_path
            
            # 清理进度显示
            progress_bar.empty()
            status_text.empty()
            
            return output_path
            
    except Exception as e:
        st.error(f"转换出错: {str(e)}")
        return None

def main():
    st.set_page_config(
        page_title="PPT转PDF工具",
        page_icon="📄",
        layout="centered"
    )
    
    st.title("PPT转PDF工具")
    st.write("上传PPT文件，自动转换为PDF格式")
    
    # 添加使用说明
    st.info("""
    使用说明：
    1. 点击"选择PPT文件"上传您的PPT文件（支持.ppt和.pptx格式）
    2. 选择PDF压缩质量
    3. 点击"转换为PDF"按钮开始转换
    4. 转换完成后，点击"下载PDF文件"保存结果
    
    注意：
    - 当前版本支持基本的文本和形状转换
    - 文件大小限制为200MB
    - 大文件处理可能需要较长时间
    - 复杂的动画效果和某些特殊格式可能无法完全保留
    """)
    
    # 添加压缩质量选择
    compression_quality = st.selectbox(
        "选择PDF压缩质量",
        ['none', 'low', 'medium', 'high'],
        format_func=lambda x: {
            'none': '不压缩',
            'low': '低质量 (文件最小)',
            'medium': '中等质量 (推荐)',
            'high': '高质量 (接近原始大小)'
        }[x]
    )
    
    uploaded_file = st.file_uploader("选择PPT文件", type=['ppt', 'pptx'])
    
    if uploaded_file is not None:
        # 检查文件大小
        file_size = len(uploaded_file.getvalue())
        if file_size > MAX_FILE_SIZE:
            st.error(f"文件太大！请上传小于 {format_size(MAX_FILE_SIZE)} 的文件（当前文件大小：{format_size(file_size)}）")
            return
            
        # 显示文件信息
        st.info(f"文件名: {uploaded_file.name}\n大小: {format_size(file_size)}")
        
        # 检查缓存
        file_hash = get_file_hash(uploaded_file.getvalue())
        cache_path = os.path.join(CACHE_DIR, f"{file_hash}_{compression_quality}.pdf")
        
        if os.path.exists(cache_path):
            st.success("找到缓存文件！")
            with open(cache_path, 'rb') as pdf_file:
                pdf_data = pdf_file.read()
                st.download_button(
                    label="下载PDF文件",
                    data=pdf_data,
                    file_name=os.path.splitext(uploaded_file.name)[0] + '.pdf',
                    mime='application/pdf'
                )
        else:
            if st.button("转换为PDF"):
                # 创建临时文件来保存上传的PPT
                with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(uploaded_file.name)[1]) as tmp_file:
                    with st.spinner('正在保存上传的文件...'):
                        progress_bar = st.progress(0)
                        save_uploaded_file(uploaded_file, tmp_file.name, progress_bar)
                        progress_bar.empty()
                    input_path = tmp_file.name
                
                with st.spinner('正在转换中...'):
                    start_time = time.time()
                    
                    # 转换文件
                    pdf_path = convert_ppt_to_pdf(input_path, compression_quality)
                    
                    if pdf_path and os.path.exists(pdf_path):
                        # 读取生成的PDF文件
                        with open(pdf_path, 'rb') as pdf_file:
                            pdf_data = pdf_file.read()
                        
                        # 保存到缓存
                        with open(cache_path, 'wb') as cache_file:
                            cache_file.write(pdf_data)
                        
                        # 获取文件大小和处理时间
                        output_size = len(pdf_data)
                        process_time = time.time() - start_time
                        
                        # 计算压缩比
                        compression_ratio = (1 - output_size / file_size) * 100
                        
                        # 提供下载链接
                        st.success(
                            f"转换成功！\n"
                            f"处理时间: {process_time:.1f}秒\n"
                            f"原始大小: {format_size(file_size)}\n"
                            f"转换后大小: {format_size(output_size)}\n"
                            f"压缩比: {compression_ratio:.1f}%"
                        )
                        st.download_button(
                            label="下载PDF文件",
                            data=pdf_data,
                            file_name=os.path.splitext(uploaded_file.name)[0] + '.pdf',
                            mime='application/pdf'
                        )
                        
                        # 清理临时文件
                        try:
                            os.remove(input_path)
                            os.remove(pdf_path)
                        except:
                            pass
                    else:
                        st.error("转换失败，请重试")

if __name__ == '__main__':
    main() 